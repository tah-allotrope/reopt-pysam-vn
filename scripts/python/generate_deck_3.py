"""Generate Deck 3 (Case 3) PPTX applying cumulative lessons from Decks 1+2.

Cumulative lessons applied (top 7 by impact):
1. Table format for Executive Summary (from Deck 1, maintained in Deck 2)
2. Slide numbers on content slides (from Deck 1, maintained in Deck 2)
3. Three section dividers (from Deck 1, maintained in Deck 2)
4. Exact hex colors locked (from Deck 1, maintained in Deck 2)
5. Card body font 11pt (from Deck 1, maintained in Deck 2)
6. Green rule line accent on every slide (from Deck 1, maintained in Deck 2)
7. Suppress footer on cover slide (NEW from Deck 2)
Bonus:
8. Add simple embedded chart on Results slide (NEW from Deck 2)
9. Full teal background on closing slide (NEW from Deck 2)
10. Partner logo placeholders (NEW from Deck 2)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors from rubric
PRIMARY_TEAL = RGBColor(0x15, 0x5B, 0x55)
TEAL_CARDS = [
    RGBColor(0x0C, 0x48, 0x3F),
    RGBColor(0x10, 0x5C, 0x51),
    RGBColor(0x11, 0x66, 0x5A),
    RGBColor(0x14, 0x76, 0x68),
    RGBColor(0x19, 0x8F, 0x7E),
]
BODY_TEXT = RGBColor(0x22, 0x22, 0x22)
ACCENT_GREEN = RGBColor(0x38, 0x76, 0x1D)
LIGHT_TEALS = [
    RGBColor(0xB5, 0xD8, 0xD0),
    RGBColor(0xD7, 0xEC, 0xE8),
    RGBColor(0xD6, 0xEC, 0xE8),
    RGBColor(0x7C, 0xD0, 0xC8),
]
NEUTRAL_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOOTER_BLACK = RGBColor(0x00, 0x00, 0x00)

def set_font(run, name="Cabin", size=12, bold=False, color=BODY_TEXT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_footer(slide, prs, slide_num, total_slides, skip_footer=False):
    """Add footer text and slide number, optionally skipping footer for cover."""
    if not skip_footer:
        footer_text = "Confidential — For Internal Use by Allotrope & Key Partners — Not Approved for Further Circulation"
        left = Inches(0.5)
        top = prs.slide_height - Inches(0.45)
        width = prs.slide_width - Inches(1.0)
        height = Inches(0.3)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = footer_text
        set_font(run, "Cabin", 8, True, FOOTER_BLACK)

    # Slide number (always add, but cover can be suppressed if desired)
    if not skip_footer:
        num_left = prs.slide_width - Inches(0.8)
        num_top = prs.slide_height - Inches(0.45)
        num_box = slide.shapes.add_textbox(num_left, num_top, Inches(0.5), Inches(0.3))
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.alignment = PP_ALIGN.RIGHT
        num_run = num_p.add_run()
        num_run.text = f"{slide_num}"
        set_font(num_run, "Cabin", 8, True, FOOTER_BLACK)

def add_green_rule(slide, prs, top_offset=Inches(0.15)):
    """Add a green accent rule line at top of slide."""
    left = Inches(0.5)
    top = top_offset
    width = Inches(1.2)
    height = Inches(0.02)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()

# ------------------------------------------------------------------
# Build deck
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

slide_num = 0

def blank_slide(skip_footer=False):
    global slide_num
    slide_num += 1
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_footer(s, prs, slide_num, 11, skip_footer=skip_footer)
    add_green_rule(s, prs)
    return s

# === Slide 1: Cover (NO FOOTER) ===
slide = blank_slide(skip_footer=True)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "DPPA Case 3\nSaigon18 Bridge Final Decision"
set_font(run, "Cabin", 26, True, PRIMARY_TEAL)

date_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(0.5))
dtf = date_box.text_frame
dp = dtf.paragraphs[0]
dp.alignment = PP_ALIGN.CENTER
drun = dp.add_run()
drun.text = "April 2026"
set_font(drun, "Cabin", 14, False, NEUTRAL_GRAY)

# Logo placeholder text
logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(2), Inches(0.4))
ltf = logo_box.text_frame
lp = ltf.paragraphs[0]
lrun = lp.add_run()
lrun.text = "ALLOTROPE"
set_font(lrun, "Cabin", 12, True, PRIMARY_TEAL)

# Partner logo placeholder
partner_box = slide.shapes.add_textbox(Inches(7.5), Inches(0.25), Inches(2), Inches(0.4))
ptf = partner_box.text_frame
pp = ptf.paragraphs[0]
pp.alignment = PP_ALIGN.RIGHT
prun = pp.add_run()
prun.text = "Partner: Tara | USAID"
set_font(prun, "Cabin", 9, False, NEUTRAL_GRAY)

# === Slide 2: Executive Summary (TABLE FORMAT) ===
slide = blank_slide()
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hp.alignment = PP_ALIGN.LEFT
hr = hp.add_run()
hr.text = "Executive Summary"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

rows, cols = 6, 2
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), Inches(9), Inches(4.0)).table

table_data = [
    ("Project", "Saigon18 DPPA Case 3"),
    ("Decision", "REJECT CURRENT CASE"),
    ("Buyer Payment", "353.9B VND/yr (vs EVN 351.0B VND/yr)"),
    ("Matched Renewable", "6.8 GWh/yr (3.7% of 184 GWh load)"),
    ("Developer DSCR", "-0.175 (required ≥ 1.0)"),
    ("Developer NPV", "-$6.05M USD (after-tax)"),
]

for i, (key, val) in enumerate(table_data):
    cell = table.cell(i, 0)
    cell.text = key
    set_font(cell.text_frame.paragraphs[0].runs[0], "Cabin", 12, True, PRIMARY_TEAL)
    cell.fill.solid()
    cell.fill.fore_color.rgb = LIGHT_TEALS[0]

    cell2 = table.cell(i, 1)
    cell2.text = val
    set_font(cell2.text_frame.paragraphs[0].runs[0], "Cabin", 12, False, BODY_TEXT)
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = WHITE

# === Slide 3: Section Divider — Context ===
slide = blank_slide()
sdiv = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.0))
sdp = sdiv.text_frame.paragraphs[0]
sdp.alignment = PP_ALIGN.CENTER
sdr = sdp.add_run()
sdr.text = "CONTEXT & APPROACH"
set_font(sdr, "Cabin", 23, True, TEAL_CARDS[0])

# === Slide 4: Key Findings (5 cards) ===
slide = blank_slide()
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hr = hp.add_run()
hr.text = "Key Findings"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

cards = [
    ("Buyer Gate", "FAIL — DPPA payment 353.9B VND/yr exceeds EVN benchmark 351.0B VND/yr by 2.9B VND/yr."),
    ("Developer Gate", "FAIL — Min DSCR = -0.175, after-tax NPV = -$6.05M. Revenue covers only 59% of debt service."),
    ("Strike Anchor", "5%-below-EVN (1,810 VND/kWh) is too low. ~3,900+ VND/kWh needed to clear developer DSCR ≥ 1.0."),
    ("BESS Impact", "Minimal matched renewable — 6.8 GWh/yr covers only 3.7% of 184 GWh annual load. 96.3% served by EVN."),
    ("Recommendation", "Reject current case. Reopen only with higher strike, smaller sizing, or actual 8760 load data."),
]

for idx, (title, body) in enumerate(cards):
    col = idx % 3
    row = idx // 3
    left = Inches(0.5 + col * 3.1)
    top = Inches(1.1 + row * 2.0)
    width = Inches(2.9)
    height = Inches(1.7)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_TEALS[idx % len(LIGHT_TEALS)]
    shape.line.color.rgb = TEAL_CARDS[idx % len(TEAL_CARDS)]
    shape.line.width = Pt(1.5)

    title_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.4))
    tp = title_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    set_font(tr, "Cabin", 14, True, TEAL_CARDS[0])

    body_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.5), width - Inches(0.2), Inches(1.1))
    bp = body_box.text_frame
    bp.word_wrap = True
    br = bp.paragraphs[0].add_run()
    br.text = body
    set_font(br, "Cabin", 11, False, BODY_TEXT)

# === Slide 5: Section Divider — Analysis ===
slide = blank_slide()
sdiv = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.0))
sdp = sdiv.text_frame.paragraphs[0]
sdp.alignment = PP_ALIGN.CENTER
sdr = sdp.add_run()
sdr.text = "ANALYSIS & RESULTS"
set_font(sdr, "Cabin", 23, True, TEAL_CARDS[0])

# === Slide 6: Financial Summary Table + Simple Chart ===
slide = blank_slide()
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hr = hp.add_run()
hr.text = "Financial Summary"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

# Left: table
rows, cols = 5, 3
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), Inches(5.5), Inches(3.5)).table
fin_data = [
    ("Metric", "Value", "Gate"),
    ("Buyer Premium", "+2.9B VND/yr", "< 0 (vs EVN)"),
    ("Matched Renewable", "6.8 GWh/yr", "—"),
    ("Developer DSCR", "-0.175", "≥ 1.0"),
    ("Developer NPV", "-$6.05M USD", "> 0"),
]
for i, (c1, c2, c3) in enumerate(fin_data):
    for j, val in enumerate([c1, c2, c3]):
        cell = table.cell(i, j)
        cell.text = val
        bold = (i == 0)
        set_font(cell.text_frame.paragraphs[0].runs[0], "Cabin", 11, bold, PRIMARY_TEAL if bold else BODY_TEXT)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_TEALS[0] if bold else WHITE

# Right: simple bar chart representation using shapes (since python-pptx can't embed native charts easily)
chart_title = slide.shapes.add_textbox(Inches(6.2), Inches(1.1), Inches(3.3), Inches(0.4))
ctp = chart_title.text_frame.paragraphs[0]
ctr = ctp.add_run()
ctr.text = "Gate Comparison"
set_font(ctr, "Cabin", 12, True, PRIMARY_TEAL)

# Bar chart bars using rectangles
bar_data = [
    ("Buyer Premium", 0.4, LIGHT_TEALS[0]),
    ("Renewable %", 0.037, LIGHT_TEALS[1]),
    ("DSCR", -0.175, RGBColor(0xFF, 0x2D, 0x78)),  # red for negative
    ("NPV", -1.0, RGBColor(0xFF, 0x2D, 0x78)),
]
bar_left = Inches(6.2)
bar_width = Inches(2.5)
for idx, (label, val, color) in enumerate(bar_data):
    top = Inches(1.6 + idx * 0.55)
    height = Inches(0.35)
    # Background bar
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, top, bar_width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()
    # Foreground bar (scaled to width)
    scaled_width = min(abs(val) * bar_width, bar_width)
    if scaled_width > Inches(0.05):
        fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, top, scaled_width, height)
        fg.fill.solid()
        fg.fill.fore_color.rgb = color
        fg.line.fill.background()
    # Label
    lbl = slide.shapes.add_textbox(bar_left + bar_width + Inches(0.1), top, Inches(1.5), height)
    lp = lbl.text_frame.paragraphs[0]
    lr = lp.add_run()
    lr.text = label
    set_font(lr, "Cabin", 9, False, BODY_TEXT)

# === Slide 7: Physical & Methodology ===
slide = blank_slide()
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hr = hp.add_run()
hr.text = "Physical Sizing & Methodology"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

bullets = [
    "Bounded-optimization REopt solve with storage floor: PV=4800kW, BESS=1500kW/3300kWh.",
    "Synthetic DPPA settlement: min(load, generation) rule with strike + adder + KPP × CFMP.",
    "Buyer benchmark: EVN bill under same TOU tariff. Developer: PySAM Single Owner model.",
    "Controller gap analysis: fixed-window dispatch vs REopt free dispatch.",
    "Site-consistent data: saigon18 load + CFMP + FMP + TOU all from same workstream.",
]
for idx, txt in enumerate(bullets):
    b = slide.shapes.add_textbox(Inches(0.7), Inches(1.1 + idx * 0.55), Inches(8.6), Inches(0.5))
    bp = b.text_frame.paragraphs[0]
    br = bp.add_run()
    br.text = "• " + txt
    set_font(br, "Cabin", 12, False, BODY_TEXT)

# === Slide 8: Section Divider — Next Steps ===
slide = blank_slide()
sdiv = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.0))
sdp = sdiv.text_frame.paragraphs[0]
sdp.alignment = PP_ALIGN.CENTER
sdr = sdp.add_run()
sdr.text = "NEXT STEPS & DECISION"
set_font(sdr, "Cabin", 23, True, TEAL_CARDS[0])

# === Slide 9: Next Steps ===
slide = blank_slide()
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hr = hp.add_run()
hr.text = "Next Steps"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

steps = [
    "Reject current case under current strike and sizing assumptions.",
    "Run strike sensitivity sweep (0%, 5%, 10%, 15%, 20% below EVN) to find viable band.",
    "Remove storage floor and run free REopt optimization for optimal PV+BESS sizing.",
    "Obtain actual monthly EVN bills and 8760 load profile for site validation.",
    "Complete 22kV two-part tariff branch analysis and compare against TOU branch.",
]
for idx, txt in enumerate(steps):
    b = slide.shapes.add_textbox(Inches(0.7), Inches(1.1 + idx * 0.55), Inches(8.6), Inches(0.5))
    bp = b.text_frame.paragraphs[0]
    br = bp.add_run()
    br.text = f"{idx+1}. " + txt
    set_font(br, "Cabin", 12, False, BODY_TEXT)

# === Slide 10: Contact / Closing (FULL TEAL BACKGROUND) ===
slide = blank_slide()
# Full teal background shape
bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = PRIMARY_TEAL
bg_shape.line.fill.background()
# Send to back
spTree = slide.shapes._spTree
sp = bg_shape._element
spTree.remove(sp)
spTree.insert(2, sp)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.8))
tp = tb.text_frame.paragraphs[0]
tp.alignment = PP_ALIGN.CENTER
tr = tp.add_run()
tr.text = "Thank You"
set_font(tr, "Cabin", 28, True, WHITE)

cb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.2))
cp = cb.text_frame.paragraphs[0]
cp.alignment = PP_ALIGN.CENTER
cr = cp.add_run()
cr.text = "Questions?\ncontact@allotrope.vc"
set_font(cr, "Cabin", 16, False, WHITE)

# Logo placeholder on closing (white text)
logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(2), Inches(0.4))
ltf = logo_box.text_frame
lp = ltf.paragraphs[0]
lrun = lp.add_run()
lrun.text = "ALLOTROPE"
set_font(lrun, "Cabin", 12, True, WHITE)

# Partner logo placeholder
partner_box = slide.shapes.add_textbox(Inches(7.5), Inches(0.25), Inches(2), Inches(0.4))
ptf = partner_box.text_frame
pp = ptf.paragraphs[0]
pp.alignment = PP_ALIGN.RIGHT
prun = pp.add_run()
prun.text = "Partner: Tara | USAID"
set_font(prun, "Cabin", 9, False, WHITE)

# Save
out_path = "reports/decks/2026-04-21-dppa-case-3.pptx"
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"Saved {out_path} ({len(prs.slides)} slides)")
