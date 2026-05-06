"""Generate Deck 2 (Case 2) PPTX applying Deck 1 lessons.

Lessons applied:
- Add slide numbers (bottom-right, Cabin 8pt)
- Use table format for Executive Summary
- Increase section dividers (Context, Analysis, Next Steps)
- Lock exact hex colors from rubric
- Match card body font size to 11-12pt
- Add green rule line accent (#38761D)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors from rubric
PRIMARY_TEAL = RGBColor(0x15, 0x5B, 0x55)
TEAL_CARDS = [
    RGBColor(0x0C, 0x48, 0x3F),
    RGBColor(0x10, 0x5C, 0x51),
    RGBColor(0x11, 0x66, 0x5A),
    RGBColor(0x14, 0x76, 0x68),
    RGBColor(0x19, 0x8F, 0x7E),
]
BODY_TEXT = RGBColor(0x22, 0x22, 0x22)
ACCENT_GREEN = RGBColor(0x38, 0x76, 0x1D)
LIGHT_TEALS = [
    RGBColor(0xB5, 0xD8, 0xD0),
    RGBColor(0xD7, 0xEC, 0xE8),
    RGBColor(0xD6, 0xEC, 0xE8),
    RGBColor(0x7C, 0xD0, 0xC8),
]
NEUTRAL_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOOTER_BLACK = RGBColor(0x00, 0x00, 0x00)

def set_font(run, name="Cabin", size=12, bold=False, color=BODY_TEXT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_footer(slide, prs, slide_num, total_slides):
    """Add footer text and slide number."""
    footer_text = "Confidential — For Internal Use by Allotrope & Key Partners — Not Approved for Further Circulation"
    left = Inches(0.5)
    top = prs.slide_height - Inches(0.45)
    width = prs.slide_width - Inches(1.0)
    height = Inches(0.3)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = footer_text
    set_font(run, "Cabin", 8, True, FOOTER_BLACK)

    # Slide number
    num_left = prs.slide_width - Inches(0.8)
    num_top = prs.slide_height - Inches(0.45)
    num_box = slide.shapes.add_textbox(num_left, num_top, Inches(0.5), Inches(0.3))
    num_tf = num_box.text_frame
    num_p = num_tf.paragraphs[0]
    num_p.alignment = PP_ALIGN.RIGHT
    num_run = num_p.add_run()
    num_run.text = f"{slide_num}"
    set_font(num_run, "Cabin", 8, True, FOOTER_BLACK)

def add_green_rule(slide, prs, top_offset=Inches(0.15)):
    """Add a green accent rule line at top of slide."""
    left = Inches(0.5)
    top = top_offset
    width = Inches(1.2)
    height = Inches(0.02)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()

# ------------------------------------------------------------------
# Build deck
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# We will manually track slide numbers
slide_num = 0

def blank_slide():
    global slide_num
    slide_num += 1
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_footer(s, prs, slide_num, 11)
    add_green_rule(s, prs)
    return s

# === Slide 1: Cover ===
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "DPPA Case 2\nNinhsim Synthetic DPPA Closeout"
set_font(run, "Cabin", 26, True, PRIMARY_TEAL)

date_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(0.5))
dtf = date_box.text_frame
dp = dtf.paragraphs[0]
dp.alignment = PP_ALIGN.CENTER
drun = dp.add_run()
drun.text = "April 2026"
set_font(drun, "Cabin", 14, False, NEUTRAL_GRAY)

# Logo placeholder text
logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(2), Inches(0.4))
ltf = logo_box.text_frame
lp = ltf.paragraphs[0]
lrun = lp.add_run()
lrun.text = "ALLOTROPE"
set_font(lrun, "Cabin", 12, True, PRIMARY_TEAL)

# === Slide 2: Executive Summary (TABLE FORMAT) ===
slide = blank_slide()
# Title
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hp.alignment = PP_ALIGN.LEFT
hr = hp.add_run()
hr.text = "Executive Summary"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

# Table: 4 rows x 2 cols
rows, cols = 5, 2
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), Inches(9), Inches(3.5)).table

table_data = [
    ("Project", "Ninhsim Synthetic DPPA (Case 2)"),
    ("Decision", "reject_current_case"),
    ("Buyer Premium", "~12.8B VND vs EVN benchmark"),
    ("Developer NPV", "Negative across all tested strikes"),
    ("Strike Range", "Tested 15%, 10%, 5%, 0% below weighted EVN"),
]

for i, (key, val) in enumerate(table_data):
    cell = table.cell(i, 0)
    cell.text = key
    set_font(cell.text_frame.paragraphs[0].runs[0], "Cabin", 12, True, PRIMARY_TEAL)
    cell.fill.solid()
    cell.fill.fore_color.rgb = LIGHT_TEALS[0]

    cell2 = table.cell(i, 1)
    cell2.text = val
    set_font(cell2.text_frame.paragraphs[0].runs[0], "Cabin", 12, False, BODY_TEXT)
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = WHITE

# === Slide 3: Section Divider — Context ===
slide = blank_slide()
sdiv = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.0))
sdp = sdiv.text_frame.paragraphs[0]
sdp.alignment = PP_ALIGN.CENTER
sdr = sdp.add_run()
sdr.text = "CONTEXT & APPROACH"
set_font(sdr, "Cabin", 23, True, TEAL_CARDS[0])

# === Slide 4: Key Findings (5 cards) ===
slide = blank_slide()
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hr = hp.add_run()
hr.text = "Key Findings"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

cards = [
    ("Buyer Premium", "~12.81B VND premium vs EVN benchmark under actual market reference."),
    ("Lowest Premium", "Even at 15% strike discount, buyer still pays ~1.55B VND premium."),
    ("Developer NPV", "Best tested developer NPV remains negative at ~-$45.19M."),
    ("CFD Stress", "Excess-generation CfD exposure adds ~4.23B VND under stress."),
    ("Recommendation", "Reject current case. Reopen only if strike basis or market data changes materially."),
]

for idx, (title, body) in enumerate(cards):
    col = idx % 3
    row = idx // 3
    left = Inches(0.5 + col * 3.1)
    top = Inches(1.1 + row * 2.0)
    width = Inches(2.9)
    height = Inches(1.7)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_TEALS[idx % len(LIGHT_TEALS)]
    shape.line.color.rgb = TEAL_CARDS[idx % len(TEAL_CARDS)]
    shape.line.width = Pt(1.5)

    title_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.4))
    tp = title_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    set_font(tr, "Cabin", 14, True, TEAL_CARDS[0])

    body_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.5), width - Inches(0.2), Inches(1.1))
    bp = body_box.text_frame
    bp.word_wrap = True
    br = bp.paragraphs[0].add_run()
    br.text = body
    set_font(br, "Cabin", 11, False, BODY_TEXT)

# === Slide 5: Section Divider — Analysis ===
slide = blank_slide()
sdiv = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.0))
sdp = sdiv.text_frame.paragraphs[0]
sdp.alignment = PP_ALIGN.CENTER
sdr = sdp.add_run()
sdr.text = "ANALYSIS & RESULTS"
set_font(sdr, "Cabin", 23, True, TEAL_CARDS[0])

# === Slide 6: Financial Summary Table ===
slide = blank_slide()
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hr = hp.add_run()
hr.text = "Financial Summary"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

rows, cols = 5, 3
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), Inches(9), Inches(3.5)).table
fin_data = [
    ("Metric", "Value", "Notes"),
    ("Buyer Premium", "12.81B VND", "vs EVN benchmark (actual market)"),
    ("Lowest Tested Premium", "1.55B VND", "At 15% strike discount"),
    ("Developer NPV (best)", "-$45.19M", "After-tax, Single Owner model"),
    ("CFD Stress Exposure", "4.23B VND", "20% market drop scenario"),
]
for i, (c1, c2, c3) in enumerate(fin_data):
    for j, val in enumerate([c1, c2, c3]):
        cell = table.cell(i, j)
        cell.text = val
        bold = (i == 0)
        set_font(cell.text_frame.paragraphs[0].runs[0], "Cabin", 12, bold, PRIMARY_TEAL if bold else BODY_TEXT)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_TEALS[0] if bold else WHITE

# === Slide 7: Settlement & Methodology ===
slide = blank_slide()
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hr = hp.add_run()
hr.text = "Settlement Design & Methodology"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

bullets = [
    "Synthetic/financial DPPA with min(load, generation) settlement quantity rule.",
    "Buyer pays (strike - market) × settled quantity; excess generation excluded.",
    "REopt.jl v0.56.4 for physical sizing; PySAM Single Owner for developer finance.",
    "Market reference replaced with repo-local saigon18 CFMP transfer series.",
    "Strike sensitivities: 15%, 10%, 5%, 0% below weighted EVN tariff.",
]
for idx, txt in enumerate(bullets):
    b = slide.shapes.add_textbox(Inches(0.7), Inches(1.1 + idx * 0.55), Inches(8.6), Inches(0.5))
    bp = b.text_frame.paragraphs[0]
    br = bp.add_run()
    br.text = "• " + txt
    set_font(br, "Cabin", 12, False, BODY_TEXT)

# === Slide 8: Section Divider — Next Steps ===
slide = blank_slide()
sdiv = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.0))
sdp = sdiv.text_frame.paragraphs[0]
sdp.alignment = PP_ALIGN.CENTER
sdr = sdp.add_run()
sdr.text = "NEXT STEPS & DECISION"
set_font(sdr, "Cabin", 23, True, TEAL_CARDS[0])

# === Slide 9: Next Steps ===
slide = blank_slide()
h = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.5))
hp = h.text_frame.paragraphs[0]
hr = hp.add_run()
hr.text = "Next Steps"
set_font(hr, "Cabin", 24, True, PRIMARY_TEAL)

steps = [
    "Keep Case 2 closed as reject under current assumptions.",
    "Reopen only if site-specific CFMP/FMP data, strike basis, or physical scope changes.",
    "If reopened, start from combined-decision JSON (Phase G) rather than replaying history.",
    "Consider hybrid settlement with excess credits as alternative mechanism.",
    "Evaluate BESS degradation modeling in future settlement revisions.",
]
for idx, txt in enumerate(steps):
    b = slide.shapes.add_textbox(Inches(0.7), Inches(1.1 + idx * 0.55), Inches(8.6), Inches(0.5))
    bp = b.text_frame.paragraphs[0]
    br = bp.add_run()
    br.text = f"{idx+1}. " + txt
    set_font(br, "Cabin", 12, False, BODY_TEXT)

# === Slide 10: Contact / Closing ===
slide = blank_slide()
tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.8))
tp = tb.text_frame.paragraphs[0]
tp.alignment = PP_ALIGN.CENTER
tr = tp.add_run()
tr.text = "Thank You"
set_font(tr, "Cabin", 28, True, PRIMARY_TEAL)

cb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.2))
cp = cb.text_frame.paragraphs[0]
cp.alignment = PP_ALIGN.CENTER
cr = cp.add_run()
cr.text = "Questions?\ncontact@allotrope.vc"
set_font(cr, "Cabin", 16, False, BODY_TEXT)

# Logo placeholder on closing
logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(2), Inches(0.4))
ltf = logo_box.text_frame
lp = ltf.paragraphs[0]
lrun = lp.add_run()
lrun.text = "ALLOTROPE"
set_font(lrun, "Cabin", 12, True, PRIMARY_TEAL)

# Save
out_path = "reports/decks/2026-04-16-dppa-case-2.pptx"
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"Saved {out_path} ({len(prs.slides)} slides)")
